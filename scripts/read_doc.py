#!/usr/bin/env python3
"""Читает .docx / .pdf / .pptx и выводит текст в stdout или сохраняет в .md.

Использование:
  python scripts/read_doc.py файл.docx
  python scripts/read_doc.py файл.pdf --output out.md
  python scripts/read_doc.py папка/ --recursive
"""

import argparse
import os
import sys
from pathlib import Path


def read_docx(path: Path) -> str:
    from docx import Document

    doc = Document(str(path))
    paras = [p.text for p in doc.paragraphs]
    # tables
    for table in doc.tables:
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells]
            paras.append(" | ".join(cells))
    return "\n".join(paras)


def read_pdf(path: Path) -> str:
    import pdfplumber

    pages = []
    with pdfplumber.open(str(path)) as pdf:
        for page in pdf.pages:
            text = page.extract_text()
            if text:
                pages.append(text)
    return "\n\n".join(pages)


def read_pptx(path: Path) -> str:
    from pptx import Presentation

    prs = Presentation(str(path))
    slides = []
    for slide in prs.slides:
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    t = para.text.strip()
                    if t:
                        texts.append(t)
        if texts:
            slides.append("\n".join(texts))
    return "\n\n---\n\n".join(slides)


def read_xlsx(path: Path) -> str:
    from openpyxl import load_workbook

    wb = load_workbook(str(path), read_only=True, data_only=True)
    parts = []
    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        rows = []
        for row in ws.iter_rows(values_only=True):
            cells = [str(c) if c is not None else "" for c in row]
            rows.append("\t".join(cells))
        if rows:
            parts.append(f"=== {sheet_name} ===\n" + "\n".join(rows))
    wb.close()
    return "\n\n".join(parts)


def read_csv(path: Path) -> str:
    import csv

    with open(str(path), newline="", encoding="utf-8-sig") as f:
        reader = csv.reader(f)
        rows = ["\t".join(row) for row in reader]
    return "\n".join(rows)


READERS = {
    ".docx": read_docx,
    ".pdf": read_pdf,
    ".pptx": read_pptx,
    ".ppt": read_pptx,
    ".xlsx": read_xlsx,
    ".xls": read_xlsx,
    ".csv": read_csv,
}


def convert_file(path: Path, output: Path | None) -> str | None:
    ext = path.suffix.lower()
    reader = READERS.get(ext)
    if reader is None:
        print(f"  [skip] неподдерживаемый формат: {ext}", file=sys.stderr)
        return None

    try:
        text = reader(path)
    except Exception as e:
        print(f"  [err] {path.name}: {e}", file=sys.stderr)
        return None

    header = f"---\nsource: {path.name}\n---\n\n"
    content = header + text

    if output:
        output.parent.mkdir(parents=True, exist_ok=True)
        output.write_text(content, encoding="utf-8")
        try:
            rel = output.relative_to(Path.cwd())
        except ValueError:
            rel = output
        print(f"  [ok]  {rel}", file=sys.stderr)
    else:
        print(content)

    return content


def main():
    parser = argparse.ArgumentParser(description="Читает .docx/.pdf/.pptx и выводит текст")
    parser.add_argument("path", help="Файл или папка")
    parser.add_argument("--output", "-o", help="Сохранить в .md (если не указан — stdout)")
    parser.add_argument("--recursive", "-r", action="store_true", help="Обработать папку рекурсивно")
    args = parser.parse_args()

    src = Path(args.path)

    if src.is_file():
        out = Path(args.output) if args.output else None
        convert_file(src, out)
        return

    if src.is_dir():
        pattern = "**/*" if args.recursive else "*"
        files = [p for p in src.glob(pattern) if p.suffix.lower() in READERS and p.is_file()]
        if not files:
            print(f"Не найдено файлов .docx/.pdf/.pptx в {src}", file=sys.stderr)
            return

        out_dir = Path(args.output) if args.output else None
        for f in sorted(files):
            out_file = None
            if out_dir:
                out_dir.mkdir(parents=True, exist_ok=True)
                out_file = out_dir / f.with_suffix(".md").name
            convert_file(f, out_file)
        return

    print(f"Путь не найден: {src}", file=sys.stderr)
    sys.exit(1)


if __name__ == "__main__":
    main()
